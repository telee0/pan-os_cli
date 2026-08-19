"""

pan-os-cli v2.5 [20260622]
pan-os-cli v2.4 [20260619]
pan-os-cli v2.3 [20260617]
pan-os-cli v2.2 [20260607]
pan-os_cli v2.1 [20260515]
pan-os_cli v2.0 [20250420]

Script to repeat CLI commands on PAN-OS over SSH

by Terence LEE <telee.hk@gmail.com>

https://github.com/telee0/pan-os_cli
https://pexpect.readthedocs.io/en/stable/index.html

"""

import json
import logging
import os
import re
import statistics
import sys
from copy import deepcopy
from datetime import datetime
from pathlib import Path

import pandas as pd
import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

from conf.cf_report import cf


ctx = {
    'start_time': datetime.now(),
}


def get_logger(name, log_file):
    logger = logging.getLogger(name)

    if logger.handlers:
        return logger

    logger.setLevel(logging.INFO)

    formatter = logging.Formatter(
        '[%(asctime)s] %(funcName)s %(levelname)s %(message)s'
    )

    file_handler = logging.FileHandler(log_file)
    file_handler.setFormatter(formatter)
    console_handler = logging.StreamHandler(sys.stdout)
    console_handler.setFormatter(formatter)

    logger.addHandler(file_handler)
    logger.addHandler(console_handler)
    logger.propagate = False

    return logger


def init():
    ctx['verbose'] = cf['verbose']
    ctx['debug'] = cf['debug']

    ctx['prs'] = Presentation(cf['template'])
    ctx['last_sldId'] = ctx['prs'].slides._sldIdLst[cf['sldId']['last']]

    ctx['cases'] = {}
    ctx['setup'] = {'bp_version': set()}
    for key in cf['pa_attrs']:
        ctx['setup'][key] = set()
        ctx['setup'][f"{key}_ex"] = set()
    ctx['results'] = []

    start_time = ctx['start_time']
    ddhhmm = start_time.strftime('%d%H%M')

    for f in('report_dir', 'log_file'):
        ctx[f] = cf[f].format(ddhhmm) if f in cf else f"{f}-{ddhhmm}"

    job_dir = ctx['report_dir']
    os.makedirs(job_dir, exist_ok=True)

    ctx['log'] = get_logger(__name__, os.path.join(job_dir, ctx['log_file']))
    ctx['log'].setLevel(logging.DEBUG if ctx['debug'] else logging.INFO if ctx['verbose'] else logging.WARNING)

    ctx['log'].info(f"initializing the environment..")

    ctx['ctx_file'] = Path(ctx['report_dir']) / cf['ctx_file']

    top_dir = Path(cf['poc_dir'])

    subdirs = sorted(
        (p for p in top_dir.iterdir() if p.is_dir()),
        key=lambda p: p.name.lower()
    )

    for prefixes in ('case_prefixes', 'other_prefixes'):
        folders = []
        for prefix in cf[prefixes]:
            prefix = prefix.casefold()
            folders.extend(
                d for d in subdirs
                if d.name.casefold().startswith(prefix)
            )
        folders_key = 'folders_' + prefixes.split('_')[0]
        ctx[folders_key] = folders
        if cf['verbose']:
            for folder in folders:
                ctx['log'].info(f"folder.name {folder.name}")


def cleanup():
    ctx['log'].info(f"cleaning up..")

    ctx['prs'].save(Path(ctx['report_dir']) / cf['report_file'])

    ctx['log'].info(f"report generation at {ctx['report_dir']} is completed, exiting..")
    ctx['end_time'] = datetime.now()
    ctx['log'].info(f"time elapsed: {ctx['end_time'] - ctx['start_time']}")

    get_joke()

    file = str(ctx['ctx_file'])
    for key, val in ctx.items():
        if isinstance(val, dict):
            for k, v in val.items():
                if isinstance(v, set):
                    val[k] = list(v)
        elif isinstance(val, list):
            ctx[key] = [v.isoformat() if isinstance(v, datetime) else v for v in val]
        # elif isinstance(val, set): ctx[key] = list(val)
        elif isinstance(val, datetime):
            ctx[key] = val.isoformat()
    with open(file, "w", encoding="utf-8") as f:
        json.dump(ctx, f, indent=2,
            skipkeys=True,
            default=str,  # lambda o: '<not serializable>',
        )


def get_joke():
    try:
        import pyjokes
        print(f"\n{pyjokes.get_joke()}")
    except Exception:
        pass


def results_to_dataframe(results):
    df = pd.DataFrame(results)

    cols = [c for c in cf['result_columns'].keys() if c in df.columns]

    df = (
        df.reindex(columns=cols)
        .sort_values(by=["case", "job"])
        .reset_index(drop=True)
    )

    df['flow_ctrl'] = df['flow_ctrl'].map(lambda x: f"{x:.0f}%")
    df['throughput'] = df['throughput'].map(lambda x: "" if x == "" else f"{x / 1000:,.1f}")
    df['flow_rate'] = df['flow_rate'].map(lambda x: "" if x == "" else f"{float(x):,.0f}")
    for col in ('connectionRate', 'allocatedSessions', 'packetRate', 'supportedSessions'):
        df[col] = df[col].map(lambda x: f"{x:,.0f}")

    df = df.rename(columns=cf['result_columns'])

    ctx['log'].info("df:\n" + df.to_string(index=False))

    return df


def slide_cover():
    slide = ctx['prs'].slides[0]
    slide.placeholders[0].text = "\n".join([cf['cust_name'], cf['subject']])
    slide.placeholders[1].text = f"{cf['poc_number']} ({cf['poc_date']})"
    slide.placeholders[2].text = cf['author']


def slide_add_bullets(text_frame, items):
    text_frame.clear()

    first = True

    def walk(items, level):
        nonlocal first

        for item in items:
            if isinstance(item, str):
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()

                p.text = item
                p.level = level

            elif isinstance(item, list):
                walk(item, level + 1)

    walk(items, 0)


def slide_add_image(slide, placeholder, image_path, scale=1.1, caption=""):
    left = placeholder.left
    top = placeholder.top
    box_w = placeholder.width
    box_h = placeholder.height + Cm(1.5)

    ctx['log'].debug(f"left {left}, top {top}, box_w {box_w}, box_h {box_h}")

    try:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
    except Exception as e:
        ctx['log'].error(f"image_path {str(image_path)}: {str(e)}")
        return

    native_w = Inches(img_w / 96) * scale
    native_h = Inches(img_h / 96) * scale

    ratio = min(
        1.0,
        box_w / native_w,
        box_h / native_h,
    )

    new_w = int(native_w * ratio)
    new_h = int(native_h * ratio)

    left += int((box_w - new_w) / 2)  # center image
    top += int((box_h - new_h) / 2)

    element = placeholder._element
    element.getparent().remove(element)  # remove placeholder

    ctx['log'].debug(
        "\n"
        + f"\timage_path '{str(image_path)}'\n"
        + f"\timg_w {img_w}, img_h {img_h}\n"
        + f"\tleft {left}, top {top}, new_w {new_w}, new_h {new_h}"
    )

    slide.shapes.add_picture(
        image_path,
        left,
        top,
        width=new_w,
        height=new_h
    )

    if caption == "":
        return

    caption_box = slide.shapes.add_textbox(
        left,
        top + new_h / 2 + Inches(0.05),
        new_w,
        Inches(0.3)
    )

    text_frame = caption_box.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    p.text = caption
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER


def slide_format(slide, replacements=None):
    if replacements is None:
        return

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key in replacements.keys():
                    placeholder = f"{{{key}}}"
                    value = " / ".join(replacements[key])
                    if placeholder in run.text:
                        run.text = run.text.replace(
                            placeholder,
                            str(value)
                        )



def slide_add_table(slide, placeholder, table_size=(2, 2)):
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height + Cm(1.5)

    ctx['log'].debug(f"left {left}, top {top}, box_w {width}, box_h {height}")

    rows, cols = table_size

    table = slide.shapes.add_table(
        rows, cols,
        left, top,
        width, height
    ).table

    table.columns[0].width = Cm(5)
    table.columns[1].width = Cm(10)

    headers = [
        "Test Case ID",
        "Description",
        "Purpose",
        "Objective",
        "Comments",
        "Metrics",
        "Result",
    ]

    data = [
        ("",), ("",), ("",), ("",), ("",),
        ("",), ("",),
    ]

    for r, text in enumerate(headers):
        table.cell(r, 0).text = text

    for r, row in enumerate(data):
        for c, value in enumerate(row, start=1):
            table.cell(r, c).text = str(value)


def slide_add_with_tables(idx, df, header=True, title=""):
    slide = ctx['prs'].slides[idx]

    shape_idx = -1
    for i, shape in enumerate(slide.shapes):
        if shape.has_table:
            shape_idx = i
            break
    if shape_idx < 0:
        return

    table = slide.shapes[shape_idx].table

    n_rows_tab = len(table.rows)
    n_cols_tab = len(table.columns)
    n_rows_df = len(df)
    n_cols_df = len(df.columns)
    n_slides = (n_rows_df + n_rows_tab - 2) // (n_rows_tab - 1)

    max_cols = min(n_cols_df, n_cols_tab)

    slide_ids = []

    i, j, n = 0, 1, n_rows_df
    while i < n:
        slide = slide_clone(idx, title=f"{title} ({j}/{n_slides})")
        slide_ids.append(list(ctx['prs'].slides).index(slide))
        table = slide.shapes[shape_idx].table

        max_rows = min(n - i, n_rows_tab - 1)

        if header:
            for c in range(max_cols):
                table.cell(0, c).text = str(df.columns[c])

        for r in range(max_rows):
            for c in range(max_cols):
                cell = table.cell(r + 1, c)
                cell.text = str(df.iat[r + i, c])

        i += max_rows
        j += 1

        for r in range(len(table.rows)):
            for c in range(len(table.columns)):
                cell = table.cell(r, c)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        run.font.name = "Montserrat"
                        run.font.size = Pt(8)
                        if r == 0:
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)

    return slide_ids


def slide_add_with_images(idx, path, images_list, title="", exclude=None):
    if images_list is None:
        return
    if exclude is None:
        exclude = []
    for images in images_list:
        if ctx['verbose']:
            ctx['log'].info(f"images {images}")
        for image_path in sorted(path.glob(images)):
            if image_path.name in exclude:
                ctx['log'].info(f"image_path {image_path} excluded")
                continue
            slide_add(idx, title=f"{title.format(image_path.stem)}", image=str(image_path))


def slide_add(idx=1, title="", text=None, image=None, caption="", table=None):
    layout = ctx['prs'].slides[idx].slide_layout
    slide = ctx['prs'].slides.add_slide(layout)
    content = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            shape.text = title
        elif shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            content = shape
            break
    if content is None:
        return slide
    if text is not None:
        content.text = text
    elif image is not None:
        slide_add_image(slide, content, image, caption=caption)
    elif table is not None:
        slide_add_table(slide, content, table_size=table)
    return slide


def slide_clone(idx=1, title=""):
    source = ctx['prs'].slides[idx]
    slide = ctx['prs'].slides.add_slide(source.slide_layout)

    for shape in list(slide.shapes):
        shape.element.getparent().remove(shape.element)

    for shape in source.shapes:
        element = deepcopy(shape.element)
        slide.shapes._spTree.insert_element_before(
            element, 'p:extLst'
        )

    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            shape.text = title
            break

    return slide


def pdf_page_get_images(doc, page, file_prefix="pdf-{}.png"):
    images = page.get_images(full=True)

    if not images:
        return []  # images is actually empty

    mask_xrefs = {image[1] for image in images if image[1] != 0}

    image_paths = []

    for i, image in enumerate(images):
        xref = image[0]

        if xref in mask_xrefs:
            continue

        pix = pymupdf.Pixmap(doc, xref)

        ctx['log'].debug(
            "\n"
            + f"xref       : {xref}"
            + f"smask      : {image[1]}"
            + f"size       : {image[2]} x {image[3]}"
            + f"bpc        : {image[4]}"
            + f"colorspace : {image[5]}"
            + f"filter     : {image[8]}"
        )
        ctx['log'].debug("\n" + "\n".join([f"image rect : {r}" for r in page.get_image_rects(xref)]))
        ctx['log'].debug(
            "\n"
            + f"pixmap     : {pix.width} x {pix.height}"
            + f"channels   : {pix.n}"
            + f"alpha      : {pix.alpha}"
        )

        image_path = os.path.join(
            ctx['report_dir'],
            f"{file_prefix.format(i)}.png"
        )

        try:
            pix.save(image_path)
            ctx['log'].info(f"image #{i} saved as '{image_path}'")
            image_paths.append(image_path)
        except Exception as e:
            ctx['log'].error(f"image not saved: {str(e)}")
            
    return image_paths


def pdf_page_find_line(page, text="", regex=None):
    blocks = page.get_text("dict")["blocks"]

    for block in blocks:
        if block["type"] != 0:
            continue
        for line in block["lines"]:
            line_text = "".join(span["text"] for span in line["spans"]).strip()
            if line_text:
                ctx['log'].debug(
                    "line_text: "
                    + f"{line['bbox'][1]:6.1f}  "
                    + f"{line['spans'][0]['size']:4.1f}  "
                    + f"'{line_text}'"
                )
            if text and text in line_text:
                return line
            if regex is not None and regex.match(line_text):
                return line

    return None


def pdf_find_section(doc, section):
    toc = doc.get_toc()

    for level, text, page in toc:
        if section.lower() in text.lower():
            return page - 1      # PyMuPDF pages are 0-based

    return -1  # not found which may cause errors


def pdf_read(pdf):
    ctx['log'].info(f"reading pdf '{pdf.name }'")

    doc = pymupdf.open(pdf)

    section_re = re.compile(cf['bp_section_re'])

    data = {}

    for i, (sect_key, section) in enumerate(cf['bp_sections'].items()):
        page_num = pdf_find_section(doc, section)
        page = doc[page_num]

        data[section] = {}

        ctx['log'].info(f"===== Section '{section}' @ Page {page_num + 1} =====")
        line = pdf_page_find_line(page, text=section)
        if line is not None:
            ctx['log'].info(f"line: text '{line['spans'][0]['text']}'")

        file_prefix = f"{pdf.stem}-{i:02d}-{page_num}-{{}}"
        image_paths = pdf_page_get_images(doc, page, file_prefix)
        data[section]['image_paths']= image_paths

        tables = page.find_tables()
        ctx['log'].info(f"section '{section}' contains {len(tables.tables)} table(s)")

        table_mode = "none"
        if sect_key in ("test_parameters", "test_device"):
            table_mode = "single"
        elif sect_key in ('super_flow_data_throughput', 'tcp_average_time_response_packet', 'transactions'):
            table_mode = "multiple"

        if table_mode == "multiple":
            rows = []
            end_of_section = False
            while page_num < len(doc):
                for table in tables.tables:
                    rows.extend(table.extract())
                page_num += 1
                page = doc[page_num]

                line = pdf_page_find_line(page, regex=section_re)
                if line is not None:
                    end_of_section = True
                    ctx['log'].info(f"end_of_section: text '{line['spans'][0]['text']}'")
                if end_of_section:
                    break

                tables = page.find_tables()
            if sect_key == "super_flow_data_throughput":
                data[section]['values'] = table_summary(rows)
            elif sect_key == "tcp_average_time_response_packet":
                data[section]['values'] = table_summary(rows, last_only=True)
            elif sect_key == "transactions":
                data[section]['values'] = table_summary(rows, last_only=True)
        elif table_mode == "single":
            heading_bottom = line['bbox'][3]
            target = None
            for table in tables.tables:
                if table.bbox[1] > heading_bottom:
                    target = table
                    break
            if target is not None:
                table = target.extract()
                if sect_key == "test_parameters":
                    data[section]['params'] = {row[0]: row[-1] for row in table[1:]}
                elif sect_key == "test_device":
                    data[section]['bp_version'] = table[1][1]
                    ctx['log'].info(f"bp_version {data[section]['bp_version']}")

    ctx['log'].info("data:\n" + "\n".join(f"\t{k}: {v}" for k, v in data.items()))

    return data


def table_summary(table, indexes=None, last_only=False):
    ctx['log'].info(f"{len(table)} rows")
    ctx['log'].debug("rows:\n" + "\n".join([f"\t{row}" for row in table]))

    if indexes is None:
        indexes = [-1]
    n_indexes = len(indexes)

    values_list = [[] for _ in indexes]

    for row in table:
        values = []
        for i in range(n_indexes):
            value = row[indexes[i]]
            try:
                values.append(float(value.replace(",", "").replace("~", "")))
            except (AttributeError, ValueError):
                pass
        if len(values) < len(indexes):
            continue
        for i, value in enumerate(values):
            values_list[i].append(value)

    summary_list = []

    if last_only:
        for values in values_list:
            summary_list.append(values[-1])
        ctx['log'].info(f"summary_list: {summary_list}")
        return summary_list

    for values in values_list:
        n = len(values)
        trim_left, trim_right = cf['bp_table_values_trim']
        left, right = int(n * trim_left), int(n * trim_right)
        values = values[left:-right] if right else values  # remove first x% and last y%
        summary = {
            'cnt': len(values),
            'min': min(values),
            'avg': statistics.mean(values),
            'med': statistics.median(values),
            'max': max(values),
        }
        summary_list.append(summary)

    ctx['log'].info(f"summary_list: {summary_list}")

    return summary_list


def pa_read(job):
    pa_info = {}

    log_file = next(job.glob(cf['cli_file']), None)
    if log_file is None:
        ctx['log'].warn("CLI log not found.")
        return {}  # pa_info

    with log_file.open("r", encoding="utf-8", errors="ignore") as f:
        lines = f.readlines()

    for attr in cf['pa_attrs']:
        for line in lines:
            if attr in line:
                value = line.partition(":")[2].strip()
                pa_info[attr] = value
                break

    ctx['log'].info("pa_info:\n" + "\n".join(f"\t{k}: {v}" for k, v in pa_info.items()))

    return pa_info


def sta_read(job):
    json_file = next(job.glob(cf['sta_file']))

    with json_file.open("r", encoding="utf-8") as f:
        data = json.load(f)

    result = {}
    for key in cf['sta_metrics'].keys():
        if key in data and "max" in data[key] and data[key]["max"]:
            result[key] = int(data[key]["max"][0])
        else:
            result[key] = None

    return result


def main():
    init()

    slide_cover()

    slide_clone(cf['sldId']['agenda'], title=cf['text']['agenda'])
    # slide = slide_add(1, title=cf['agenda'])
    # slide_add_bullets(slide.placeholders[1].text_frame, cf['agenda_items'])

    for folder in ctx['folders_case']:
        case = folder.name
        print(". " * 40)
        ctx['log'].info(f"case '{case}'")

        slide_clone(cf['sldId']['section'], title=case)  # slide for section start
        slide_clone(cf['sldId']['case'], title=f"{cf['text']['case']} - {case}")  # slide for test case details

        ctx['cases'][case] = {}
        throughput, flow_rate = "", ""
        pdf_list = list(folder.glob(cf['bp_reports']))
        for pdf in pdf_list:
            ctx['log'].info(f"case '{case}' pdf '{pdf.name}'")
            data = pdf_read(pdf)
            ctx['cases'][case][pdf.stem] = data
            ctx['setup']['bp_version'].add(data[cf['bp_sections']['test_device']]['bp_version'])
            throughput = data[cf['bp_sections']['super_flow_data_throughput']]['values'][-1]['avg']
            if 'flow_rate' not in ctx['cases'][case]:
                ctx['cases'][case]['flow_rate'] = []
            flow_rate = data[cf['bp_sections']['test_parameters']]['params']['Maximum Flow Creation Rate']
            ctx['cases'][case]['flow_rate'].append(flow_rate)
        for pdf in ctx['cases'][case]:
            data = ctx['cases'][case][pdf]
            for section in cf['bp_sections'].values():
                if section not in data:
                    continue
                caption = str(data[section]['values']) if 'values' in data[section] else ""
                for image in data[section]['image_paths']:
                    slide_add(cf['sldId']['new'], title=f"{case} - {section}", image=image, caption=caption)

        job_list = list(folder.glob(cf['job_dir']))
        for job in job_list:
            ctx['cases'][case][job.stem] = {}
            ctx['log'].info(f"case '{case}' job '{job.name}'")
            pa_info = pa_read(job)
            ctx['cases'][case][job.stem]['pa'] = pa_info

            for key in cf['pa_attrs']:
                ctx['setup'][key].add(pa_info[key])
            ctx['setup']['sw-version_ex'].add(f"{pa_info['sw-version']} ({pa_info['model']})")
            ctx['setup']['app-version_ex'].add(f"{pa_info['app-version']} ({pa_info['app-release-date'][:10]})")

            stats = sta_read(job)
            stats.update({
                'case': case,
                'job': job.name,
                'throughput': throughput,
                'flow_rate': flow_rate,
                'supportedSessions': int(pa_info['sessions supported'])
            })
            throughput, flow_rate = "", ""  # used once
            ctx['results'].append(stats)
            slide_add_with_images(cf['sldId']['new'], job, cf['dp_files_list'], title=f"{case} - {cf['text']['dp']} ({{}})")
            slide_add_with_images(cf['sldId']['new'], job, cf['p_files_list'], title=f"{case} - {cf['text']['util']} ({{}})")

        slide_add_with_images(cf['sldId']['new'], folder, [cf['other_files']], exclude=cf['other_files_exclude'], title=f"{case} - {{}}")

    slide_clone(cf['sldId']['section'], title=cf['text']['others'])  # slide for section start
    for folder in ctx['folders_other']:
        slide_add_with_images(cf['sldId']['new'], folder, [cf['image_files']], title=f"{cf['text']['others']} - {folder.name} ({{}})")

    slide = ctx['prs'].slides[cf['sldId']['setup']]
    slide_format(slide, ctx['setup'])

    df = results_to_dataframe(ctx['results'])
    idx_list = slide_add_with_tables(cf['sldId']['results'], df, header=True, title=cf['text']['results'])
    print(f"idx_list: {idx_list}")
    slides = ctx['prs'].slides._sldIdLst
    for i, idx in enumerate(idx_list):
        sldId = slides[idx]
        slides.remove(sldId)
        slides.insert(cf['sldId']['results'] + i, sldId)  # move the result slides to the target position

    slides = ctx['prs'].slides._sldIdLst
    slides.remove(ctx['last_sldId'])
    slides.append(ctx['last_sldId'])  # move the previous last slide to the end of the deck

    if 'removal' in cf['sldId']:
        left, right = cf['sldId']['removal']
        slides_removal = list(slides[left:right])
        for sldId in slides_removal:
            slides.remove(sldId)

    cleanup()

if __name__ == '__main__':
    main()
